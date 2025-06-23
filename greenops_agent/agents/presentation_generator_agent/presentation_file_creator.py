from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from io import BytesIO
import requests
from lxml import etree
from google.adk.tools import ToolContext
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload
from google.oauth2 import service_account
import json
import os
from greenops_agent.secrets_access_manager import access_secret

def get_shape_by_name(slide, target_name):
    for shape in slide.shapes:
        if shape.name == target_name:
            return shape
    return None

def set_text_with_optional_style(shape, text, font_size=None, font_color=None, bold=None):
    if not shape.has_text_frame:
        return

    text_frame = shape.text_frame
    text_frame.clear()  # Optional: clear existing text if you want full control

    for i, line in enumerate(text.split("\n")):
        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        run = p.add_run()

        line = line.strip()

        if line.startswith("-"):
            run.text = line[1:].strip()
        else:
            run.text = line
            p._pPr.insert(0, etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"))

        if text=="THANK YOU!":
            p._pPr.insert(0, etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"))

        font = run.font
        if font_size:
            font.size = font_size
        if font_color:
            font.color.rgb = font_color
        if bold is not None:
            font.bold = bold

def create_presentation(content: dict, tool_context: ToolContext):
    """
    Input: Dict containting the data for slides
    """
    CHART_IMAGE_MAP = tool_context.state["chart_links"]
    content["forecast_overview"]["chart_image_uri"] = CHART_IMAGE_MAP["[[chart_carbon_timeseries]]"]
    content["regional_utilization"]["chart_image_uri"] = CHART_IMAGE_MAP["[[chart_underutilization]]"]
    content["top_recommendations"]["chart_image_uri"] = CHART_IMAGE_MAP["[[chart_region_utilization]]"]
    content["instance_behavior_insights"]["chart_image_uri"] = CHART_IMAGE_MAP["[[chart_cpu_vs_carbon]]"]

    prs = Presentation("custom_template.pptx")

    prs.slides._sldIdLst.remove(list(prs.slides._sldIdLst)[0])

    # HERO PAGE
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    set_text_with_optional_style(get_shape_by_name(slide1, "Title 1"), "GreenOps Weekly Report", font_size=Pt(50))
    set_text_with_optional_style(get_shape_by_name(slide1, "Subtitle 2"), content["hero_page"]["week_date_range"], font_size=Pt(20))

    # EXECUTIVE SUMMARY
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    set_text_with_optional_style(get_shape_by_name(slide2, "Title 1"), "Executive Summary", font_size=Pt(30))
    set_text_with_optional_style(get_shape_by_name(slide2, "Text Placeholder 2"), content["executive_summary"]["content"], font_size=Pt(28))


    # FORECAST OVERVIEW
    slide3 = prs.slides.add_slide(prs.slide_layouts[2])
    set_text_with_optional_style(get_shape_by_name(slide3, "Title 1"), "Forecast Overview", font_size=Pt(30))

    image_shape = get_shape_by_name(slide3, "Picture Placeholder 2")
    image_url = content["forecast_overview"]["chart_image_uri"]

    if image_shape:
        response = requests.get(image_url)
        if response.status_code == 200:
            image_data = BytesIO(response.content)
            picture = image_shape.insert_picture(image_data)
            picture.crop_top = 0
            picture.crop_left = 0
            picture.crop_bottom = 0
            picture.crop_right = 0


    set_text_with_optional_style(get_shape_by_name(slide3, "Text Placeholder 3"), content["forecast_overview"]["content"])

    # Regional Utilization
    slide4 = prs.slides.add_slide(prs.slide_layouts[3])
    set_text_with_optional_style(get_shape_by_name(slide4, "Title 1"), "Regional Utilization", font_size=Pt(30))

    image_shape = get_shape_by_name(slide4, "Picture Placeholder 2")
    image_url = content["regional_utilization"]["chart_image_uri"]

    if image_shape:
        response = requests.get(image_url)
        if response.status_code == 200:
            image_data = BytesIO(response.content)
            picture = image_shape.insert_picture(image_data)
            picture.crop_top = 0
            picture.crop_left = 0
            picture.crop_bottom = 0
            picture.crop_right = 0


    set_text_with_optional_style(get_shape_by_name(slide4, "Text Placeholder 3"), content["regional_utilization"]["content"])


    # Top Recommendations
    slide5 = prs.slides.add_slide(prs.slide_layouts[4])
    set_text_with_optional_style(get_shape_by_name(slide5, "Title 1"), "Top Recommendations", font_size=Pt(30))

    image_shape = get_shape_by_name(slide5, "Picture Placeholder 2")
    image_url = content["top_recommendations"]["chart_image_uri"]

    if image_shape:
        response = requests.get(image_url)
        if response.status_code == 200:
            image_data = BytesIO(response.content)
            picture = image_shape.insert_picture(image_data)
            picture.crop_top = 0
            picture.crop_left = 0
            picture.crop_bottom = 0
            picture.crop_right = 0


    set_text_with_optional_style(get_shape_by_name(slide5, "Text Placeholder 3"), content["top_recommendations"]["content"], font_size=Pt(20))


    # Top Recommendations
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    set_text_with_optional_style(get_shape_by_name(slide6, "Title 1"), "Instance Behaviour Insights", font_size=Pt(30))

    image_shape = get_shape_by_name(slide6, "Picture Placeholder 2")
    image_url = content["instance_behavior_insights"]["chart_image_uri"]

    if image_shape:
        response = requests.get(image_url)
        if response.status_code == 200:
            image_data = BytesIO(response.content)
            picture = image_shape.insert_picture(image_data)
            picture.crop_top = 0
            picture.crop_left = 0
            picture.crop_bottom = 0
            picture.crop_right = 0


    set_text_with_optional_style(get_shape_by_name(slide6, "Text Placeholder 3"), content["instance_behavior_insights"]["content"], font_size=Pt(20))


    # Thank you!
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_text_with_optional_style(get_shape_by_name(slide7, "Text Placeholder 1"), "THANK YOU!", font_size=Pt(70))

    if not os.path.exists("presentations/"):
        os.mkdir("presentations/")

    filename = f"GreenOps Weekly Report {content['hero_page']['week_date_range']}.pptx"
    file_path = f"presentations/{filename}"
    prs.save(file_path)

    drive_link = upload_pptx_and_get_download_link(file_path, filename)

    os.remove(file_path)

    return {
        "Download_link": drive_link
    }


def upload_pptx_and_get_download_link(filepath, filename):
    SCOPES = ['https://www.googleapis.com/auth/drive']

    if not os.environ.get("SERVICE_ACCOUNT_KEY"):
        os.environ["SERVICE_ACCOUNT_KEY"] = access_secret(secret_id="SERVICE_ACCOUNT_KEY")

    creds = service_account.Credentials.from_service_account_info(json.loads(os.environ["SERVICE_ACCOUNT_KEY"]), scopes=SCOPES)

    drive_service = build('drive', 'v3', credentials=creds)

    media = MediaFileUpload(filepath, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

    file_metadata = {
        'name': filename,
        'mimeType': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    }

    uploaded_file = drive_service.files().create(
        body=file_metadata,
        media_body=media,
        fields='id'
    ).execute()

    file_id = uploaded_file['id']

    # Make it public (optional)
    drive_service.permissions().create(
        fileId=file_id,
        body={'type': 'anyone', 'role': 'reader'},
    ).execute()

    # Download link (not preview)
    return f"https://drive.google.com/uc?id={file_id}&export=download"


print(upload_pptx_and_get_download_link("/Users/nikhilmankani/Documents/GreenOps/GreenOps/output1.pptx", "GreenOpsFirstPresentation.pptx"))